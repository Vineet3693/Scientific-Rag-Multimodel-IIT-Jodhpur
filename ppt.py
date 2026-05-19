
# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

# ── Color Palette ──────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # Dark Navy
ACCENT_CYAN  = RGBColor(0x00, 0xB4, 0xD8)   # Cyan
ACCENT_GREEN = RGBColor(0x06, 0xD6, 0xA0)   # Green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)   # White
LIGHT_GRAY   = RGBColor(0xAB, 0xB2, 0xBF)   # Gray
CARD_BG      = RGBColor(0x1B, 0x2A, 0x3B)   # Card Dark Blue

# ── Helper Functions ────────────────────────────────────
def add_bg(slide, prs, color=DARK_BG):
    """Add solid background color to slide"""
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    """Add a styled textbox"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_card(slide, left, top, width, height, color=CARD_BG):
    """Add a card/box background"""
    card = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = ACCENT_CYAN
    card.line.width = Pt(1)
    return card

def add_divider(slide, top, prs):
    """Add horizontal cyan divider line"""
    line = slide.shapes.add_shape(
        1,
        Inches(0.5), Inches(top),
        Inches(12.3), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.fill.background()

def add_multiline_textbox(slide, lines, left, top,
                           width, height,
                           font_size=14, color=WHITE):
    """Add textbox with multiple styled lines"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, fsize, fcolor) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fsize or font_size)
        run.font.bold = bold
        run.font.color.rgb = fcolor or color

# ── Create Presentation ─────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
add_bg(slide1, prs)

# Accent top bar
bar = slide1.shapes.add_shape(
    1, 0, 0,
    prs.slide_width, Inches(0.08)
)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_CYAN
bar.line.fill.background()

# Main Title
add_textbox(slide1,
    "RAG Pipeline Architecture",
    1.0, 1.8, 11.0, 1.2,
    font_size=44, bold=True,
    color=ACCENT_CYAN,
    align=PP_ALIGN.CENTER
)

# Subtitle
add_textbox(slide1,
    "Retrieval Augmented Generation – From Data to Intelligent Response",
    1.0, 3.0, 11.0, 0.6,
    font_size=20, bold=False,
    color=WHITE,
    align=PP_ALIGN.CENTER
)

add_divider(slide1, 3.85, prs)

# Author
add_textbox(slide1,
    "Vineet Yadav  |  AI Engineer  |  LLM Orchestration Specialist",
    1.0, 4.1, 11.0, 0.5,
    font_size=16, bold=False,
    color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER
)

# Tech tags
tags = ["LangChain", "ChromaDB", "FAISS", "Gemini", "Groq", "Streamlit"]
for i, tag in enumerate(tags):
    x = 1.5 + i * 1.72
    tag_box = slide1.shapes.add_shape(
        1,
        Inches(x), Inches(5.2),
        Inches(1.4), Inches(0.4)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = CARD_BG
    tag_box.line.color.rgb = ACCENT_CYAN
    tag_box.line.width = Pt(1)
    add_textbox(slide1, tag,
        x + 0.05, 5.22,
        1.3, 0.35,
        font_size=11, bold=True,
        color=ACCENT_CYAN,
        align=PP_ALIGN.CENTER
    )

# Bottom bar
bot = slide1.shapes.add_shape(
    1, 0, Inches(7.3),
    prs.slide_width, Inches(0.2)
)
bot.fill.solid()
bot.fill.fore_color.rgb = CARD_BG
bot.line.fill.background()

# ════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_bg(slide2, prs)

add_textbox(slide2,
    "Why Do LLMs Hallucinate?",
    0.5, 0.3, 10.0, 0.8,
    font_size=32, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide2, 1.2, prs)

problems = [
    "❌  LLMs have knowledge cutoff dates",
    "❌  No access to private / domain-specific data",
    "❌  Cannot cite sources accurately",
    "❌  Expensive to fine-tune for every use case",
    "❌  Responses not grounded in real documents",
]
for i, prob in enumerate(problems):
    add_card(slide2, 0.5, 1.5 + i * 0.9, 7.5, 0.75)
    add_textbox(slide2, prob,
        0.7, 1.55 + i * 0.9,
        7.2, 0.65,
        font_size=16, color=WHITE
    )

# Solution card
add_card(slide2, 8.5, 1.5, 4.3, 4.5,
         color=RGBColor(0x06, 0x3A, 0x2E))
add_textbox(slide2, "✅  Solution",
    8.7, 1.7, 4.0, 0.5,
    font_size=18, bold=True,
    color=ACCENT_GREEN
)
add_textbox(slide2,
    "RAG Pipeline\nRetrieval Augmented\nGeneration",
    8.7, 2.3, 4.0, 1.5,
    font_size=22, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER
)
add_textbox(slide2,
    "Ground LLM responses\nin real documents\nand live data",
    8.7, 3.9, 4.0, 1.2,
    font_size=14, color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER
)

# ════════════════════════════════════════════════════════
# SLIDE 3 — What is RAG?
# ════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_bg(slide3, prs)

add_textbox(slide3,
    "What is RAG?",
    0.5, 0.3, 10.0, 0.8,
    font_size=32, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide3, 1.2, prs)

add_textbox(slide3,
    "RAG = Retrieval System  +  LLM Generation",
    0.5, 1.4, 12.0, 0.7,
    font_size=22, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER
)

components = [
    ("📥", "Knowledge Base",
     "Your private documents\nPDFs, DOCX, TXT, Web",
     RGBColor(0x1B, 0x3A, 0x5C)),
    ("🔍", "Retrieval Engine",
     "Finds relevant context\nSemantic + Hybrid Search",
     RGBColor(0x1B, 0x3A, 0x3A)),
    ("🤖", "LLM Generator",
     "Produces accurate response\nGrounded in real data",
     RGBColor(0x2A, 0x1B, 0x3A)),
]
for i, (icon, title, desc, bg_col) in enumerate(components):
    x = 0.6 + i * 4.2
    add_card(slide3, x, 2.3, 3.8, 3.5, color=bg_col)
    add_textbox(slide3, icon,
        x, 2.5, 3.8, 0.8,
        font_size=36, align=PP_ALIGN.CENTER
    )
    add_textbox(slide3, title,
        x, 3.3, 3.8, 0.6,
        font_size=18, bold=True,
        color=ACCENT_CYAN,
        align=PP_ALIGN.CENTER
    )
    add_textbox(slide3, desc,
        x + 0.2, 3.95, 3.4, 1.5,
        font_size=13, color=LIGHT_GRAY,
        align=PP_ALIGN.CENTER
    )

add_textbox(slide3,
    '"Give the LLM the right context at the right time"',
    0.5, 6.3, 12.3, 0.6,
    font_size=15, italic=True,
    color=ACCENT_GREEN,
    align=PP_ALIGN.CENTER
)

# ════════════════════════════════════════════════════════
# SLIDE 4 — Full Pipeline Architecture
# ════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_bg(slide4, prs)

add_textbox(slide4,
    "End-to-End RAG Pipeline Flow",
    0.5, 0.3, 12.0, 0.7,
    font_size=32, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide4, 1.1, prs)

steps = [
    ("📄 Raw Documents",   "PDF / DOCX / TXT"),
    ("✂️ Text Chunking",    "512–1024 token splits"),
    ("🔢 Embedding Model", "Text → Vectors"),
    ("🗄️ Vector Database", "ChromaDB / FAISS"),
    ("🔍 Similarity Search","Top-K Retrieval"),
    ("📝 Prompt Template", "Context + Query"),
    ("🤖 LLM Generation",  "Gemini / GPT / Groq"),
    ("✅ Final Response",   "Accurate & Cited"),
]
box_w = 1.35
box_h = 0.75
start_x = 0.3
start_y = 1.5

for i, (step, desc) in enumerate(steps):
    x = start_x + i * 1.58
    add_card(slide4, x, start_y, box_w, box_h)
    add_textbox(slide4, step,
        x + 0.05, start_y + 0.02,
        box_w - 0.1, 0.38,
        font_size=10, bold=True,
        color=ACCENT_CYAN,
        align=PP_ALIGN.CENTER
    )
    add_textbox(slide4, desc,
        x + 0.05, start_y + 0.38,
        box_w - 0.1, 0.3,
        font_size=8, color=LIGHT_GRAY,
        align=PP_ALIGN.CENTER
    )
    # Arrow
    if i < len(steps) - 1:
        add_textbox(slide4, "→",
            x + box_w, start_y + 0.2,
            0.2, 0.4,
            font_size=14, bold=True,
            color=ACCENT_CYAN,
            align=PP_ALIGN.CENTER
        )

# Two phase labels
add_card(slide4, 0.3, 2.5, 5.8, 0.5,
         color=RGBColor(0x0A, 0x25, 0x3A))
add_textbox(slide4, "⬆️  INDEXING PHASE (Offline)",
    0.4, 2.55, 5.6, 0.4,
    font_size=12, bold=True,
    color=ACCENT_CYAN,
    align=PP_ALIGN.CENTER
)

add_card(slide4, 6.3, 2.5, 6.7, 0.5,
         color=RGBColor(0x0A, 0x2A, 0x25))
add_textbox(slide4, "⬆️  QUERYING PHASE (Online)",
    6.4, 2.55, 6.5, 0.4,
    font_size=12, bold=True,
    color=ACCENT_GREEN,
    align=PP_ALIGN.CENTER
)

# Architecture summary box
add_card(slide4, 0.3, 3.2, 12.7, 3.8,
         color=RGBColor(0x10, 0x20, 0x30))
add_textbox(slide4,
    "Complete RAG Architecture",
    0.5, 3.3, 12.3, 0.5,
    font_size=16, bold=True,
    color=ACCENT_CYAN,
    align=PP_ALIGN.CENTER
)

architecture_text = [
    ("📥  INPUT LAYER       →  PDF Loader | DOCX Parser | Web Scraper | Text Files", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("✂️   PROCESSING LAYER  →  Text Splitter | Chunk Overlap | Token Counter", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("🔢  EMBEDDING LAYER   →  OpenAI Embeddings | HuggingFace | Sentence Transformers", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("🗄️   STORAGE LAYER     →  ChromaDB | FAISS | Pinecone | Weaviate", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("🔍  RETRIEVAL LAYER   →  Semantic Search | BM25 | Hybrid Search | Reranking", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("🤖  GENERATION LAYER  →  GPT-4 | Gemini Pro | Groq LLaMA | Claude", False, 13, WHITE),
]
add_multiline_textbox(slide4, architecture_text,
    0.5, 3.85, 12.3, 3.0)

# ════════════════════════════════════════════════════════
# SLIDE 5 — Data Ingestion Layer
# ════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
add_bg(slide5, prs)

add_textbox(slide5,
    "Step 1 – Document Processing Pipeline",
    0.5, 0.3, 12.0, 0.7,
    font_size=32, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide5, 1.1, prs)

# Table Headers
headers = ["Stage", "Process", "Tools Used", "Output"]
col_widths = [1.8, 3.0, 3.5, 4.5]
col_positions = [0.4, 2.3, 5.4, 8.9]

# Header row
for j, (header, cw, cx) in enumerate(
        zip(headers, col_widths, col_positions)):
    header_card = slide5.shapes.add_shape(
        1,
        Inches(cx), Inches(1.3),
        Inches(cw - 0.1), Inches(0.5)
    )
    header_card.fill.solid()
    header_card.fill.fore_color.rgb = ACCENT_CYAN
    header_card.line.fill.background()
    add_textbox(slide5, header,
        cx + 0.05, 1.32,
        cw - 0.2, 0.42,
        font_size=13, bold=True,
        color=DARK_BG,
        align=PP_ALIGN.CENTER
    )

# Table rows
table_data = [
    ("📥 Input",     "Load Documents",          "LangChain Loaders",          "Raw Text Extracted"),
    ("✂️ Chunking",  "Split into 512-1024 tokens","RecursiveTextSplitter",     "Text Chunks Created"),
    ("🧹 Cleaning",  "Remove noise & format",    "Python, Regex, NLTK",        "Clean Text Ready"),
    ("🔢 Embedding", "Convert text → vectors",   "OpenAI / HuggingFace",       "Float Vectors Generated"),
    ("💾 Storage",   "Store in vector DB",        "ChromaDB / FAISS",           "Indexed & Searchable"),
]

for i, (stage, process, tools, output) in enumerate(table_data):
    row_color = CARD_BG if i % 2 == 0 else RGBColor(0x14, 0x24, 0x35)
    row_data = [stage, process, tools, output]
    for j, (cell, cw, cx) in enumerate(
            zip(row_data, col_widths, col_positions)):
        cell_box = slide5.shapes.add_shape(
            1,
            Inches(cx), Inches(1.9 + i * 0.75),
            Inches(cw - 0.1), Inches(0.65)
        )
        cell_box.fill.solid()
        cell_box.fill.fore_color.rgb = row_color
        cell_box.line.color.rgb = RGBColor(0x2A, 0x3A, 0x4A)
        cell_box.line.width = Pt(0.5)
        add_textbox(slide5, cell,
            cx + 0.1, 1.92 + i * 0.75,
            cw - 0.2, 0.55,
            font_size=12, color=WHITE
        )

# Project highlight box
add_card(slide5, 0.4, 5.8, 12.5, 1.4,
         color=RGBColor(0x06, 0x2A, 0x1E))
add_textbox(slide5,
    "🚀  Your Project Result:",
    0.6, 5.9, 4.0, 0.4,
    font_size=14, bold=True,
    color=ACCENT_GREEN
)
add_textbox(slide5,
    "6 Chemical Engineering Textbooks  →  1,666 Chunks  →  ChromaDB  →  ~90% Answer Relevance ✅",
    0.6, 6.35, 12.2, 0.5,
    font_size=14, bold=False,
    color=WHITE
)

# ════════════════════════════════════════════════════════
# SLIDE 6 — Retrieval Engine Layer
# ════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
add_bg(slide6, prs)

add_textbox(slide6,
    "Step 2 – Hybrid Retrieval System",
    0.5, 0.3, 12.0, 0.7,
    font_size=32, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide6, 1.1, prs)

# Semantic Search Card
add_card(slide6, 0.4, 1.3, 5.8, 4.5,
         color=RGBColor(0x0A, 0x20, 0x3A))
add_textbox(slide6, "🔵  Semantic Search",
    0.5, 1.45, 5.6, 0.55,
    font_size=20, bold=True,
    color=ACCENT_CYAN
)
semantic_points = [
    "→  Converts query to embedding vector",
    "→  Finds similar vectors (cosine similarity)",
    "→  Understands meaning not just keywords",
    "→  Best for: Conceptual questions",
    "→  Uses: ChromaDB / FAISS index",
]
for i, point in enumerate(semantic_points):
    add_textbox(slide6, point,
        0.6, 2.1 + i * 0.55, 5.5, 0.5,
        font_size=13, color=WHITE
    )

# Internet Search Card
add_card(slide6, 6.8, 1.3, 5.8, 4.5,
         color=RGBColor(0x0A, 0x2A, 0x1A))
add_textbox(slide6, "🟢  Internet / Keyword Search",
    6.9, 1.45, 5.6, 0.55,
    font_size=20, bold=True,
    color=ACCENT_GREEN
)
internet_points = [
    "→  Real-time web search integration",
    "→  Fetches current / live information",
    "→  Best for: Recent events & live data",
    "→  Uses: SerpAPI / DuckDuckGo / Tavily",
    "→  Combines with semantic for accuracy",
]
for i, point in enumerate(internet_points):
    add_textbox(slide6, point,
        6.95, 2.1 + i * 0.55, 5.5, 0.5,
        font_size=13, color=WHITE
    )

# Hybrid Center
add_textbox(slide6, "⭐",
    6.0, 2.8, 0.8, 0.8,
    font_size=28, align=PP_ALIGN.CENTER
)
add_textbox(slide6, "+",
    6.1, 3.5, 0.6, 0.5,
    font_size=28, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER
)

# Result bar
add_card(slide6, 0.4, 6.0, 12.5, 1.1,
         color=RGBColor(0x1A, 0x35, 0x25))
add_textbox(slide6,
    "⭐  Hybrid = Semantic + Internet Search  →  ~90% Improvement in Answer Relevance ✅",
    0.6, 6.15, 12.0, 0.6,
    font_size=15, bold=True,
    color=ACCENT_GREEN,
    align=PP_ALIGN.CENTER
)

# ════════════════════════════════════════════════════════
# SLIDE 7 — Generation Layer
# ════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank_layout)
add_bg(slide7, prs)

add_textbox(slide7,
    "Step 3 – LLM Response Generation",
    0.5, 0.3, 12.0, 0.7,
    font_size=32, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide7, 1.1, prs)

# Prompt Template Box
add_card(slide7, 0.4, 1.3, 6.5, 4.8,
         color=RGBColor(0x10, 0x18, 0x28))
add_textbox(slide7, "📝  Prompt Template Structure",
    0.6, 1.4, 6.2, 0.5,
    font_size=16, bold=True,
    color=ACCENT_CYAN
)

prompt_lines = [
    ("SYSTEM:", True, 12, ACCENT_CYAN),
    ("  You are an expert assistant.", False, 11, LIGHT_GRAY),
    ("  Answer ONLY using context below.", False, 11, LIGHT_GRAY),
    ("", False, 6, WHITE),
    ("CONTEXT:", True, 12, ACCENT_CYAN),
    ("  {retrieved_chunks}", False, 11, ACCENT_GREEN),
    ("", False, 6, WHITE),
    ("QUESTION:", True, 12, ACCENT_CYAN),
    ("  {user_query}", False, 11, ACCENT_GREEN),
    ("", False, 6, WHITE),
    ("ANSWER:", True, 12, ACCENT_CYAN),
    ("  [LLM generates grounded response]", False, 11, LIGHT_GRAY),
]
add_multiline_textbox(slide7, prompt_lines,
    0.6, 2.0, 6.1, 3.8)

# LLM Comparison Table
add_textbox(slide7, "🤖  LLM Comparison",
    7.2, 1.3, 5.8, 0.5,
    font_size=16, bold=True,
    color=ACCENT_CYAN
)

llm_headers = ["Model", "Speed", "Cost", "Quality"]
llm_col_w   = [1.5, 1.1, 1.0, 1.9]
llm_col_x   = [7.2, 8.75, 9.9, 10.95]

for j, (h, cw, cx) in enumerate(
        zip(llm_headers, llm_col_w, llm_col_x)):
    hbox = slide7.shapes.add_shape(
        1,
        Inches(cx), Inches(1.9),
        Inches(cw - 0.05), Inches(0.45)
    )
    hbox.fill.solid()
    hbox.fill.fore_color.rgb = ACCENT_CYAN
    hbox.line.fill.background()
    add_textbox(slide7, h,
        cx + 0.05, 1.92,
        cw - 0.1, 0.38,
        font_size=11, bold=True,
        color=DARK_BG,
        align=PP_ALIGN.CENTER
    )

llm_data = [
    ("GPT-4",       "Medium",    "$$$$", "⭐⭐⭐⭐⭐"),
    ("Gemini Pro",  "Fast",      "$$$",  "⭐⭐⭐⭐"),
    ("Groq LLaMA",  "Very Fast", "$",    "⭐⭐⭐⭐"),
    ("Claude 3",    "Fast",      "$$$",  "⭐⭐⭐⭐⭐"),
    ("Mistral",     "Fast",      "$$",   "⭐⭐⭐"),
]
for i, row in enumerate(llm_data):
    rc = CARD_BG if i % 2 == 0 else RGBColor(0x14, 0x24, 0x35)
    for j, (cell, cw, cx) in enumerate(
            zip(row, llm_col_w, llm_col_x)):
        cb = slide7.shapes.add_shape(
            1,
            Inches(cx), Inches(2.45 + i * 0.65),
            Inches(cw - 0.05), Inches(0.58)
        )
        cb.fill.solid()
        cb.fill.fore_color.rgb = rc
        cb.line.color.rgb = RGBColor(0x2A, 0x3A, 0x4A)
        cb.line.width = Pt(0.5)
        add_textbox(slide7, cell,
            cx + 0.05, 2.47 + i * 0.65,
            cw - 0.1, 0.5,
            font_size=11, color=WHITE,
            align=PP_ALIGN.CENTER
        )

# ════════════════════════════════════════════════════════
# SLIDE 8 — Real Project Results
# ════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
add_bg(slide8, prs)

add_textbox(slide8,
    "Chemical Engineering RAG System – Live Results",
    0.5, 0.3, 12.0, 0.7,
    font_size=30, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide8, 1.1, prs)

# Metrics Cards Row
metrics = [
    ("📊", "Answer\nRelevance",  "~45%",  "→",  "~90%",  ACCENT_GREEN),
    ("🤝", "Source\nCitation",   "❌ None","→",  "✅ PDF", ACCENT_GREEN),
    ("🧠", "Hallucination\nRate","High",   "→",  "Very Low",ACCENT_CYAN),
    ("🎯", "Domain\nAccuracy",   "Low",    "→",  "High",  ACCENT_CYAN),
]
for i, (icon, label, before, arrow, after, color) in enumerate(metrics):
    x = 0.4 + i * 3.15
    add_card(slide8, x, 1.3, 2.9, 2.8)
    add_textbox(slide8, icon,
        x, 1.4, 2.9, 0.6,
        font_size=28, align=PP_ALIGN.CENTER
    )
    add_textbox(slide8, label,
        x, 2.0, 2.9, 0.6,
        font_size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER
    )
    add_textbox(slide8, before,
        x + 0.1, 2.65, 0.9, 0.45,
        font_size=12, color=LIGHT_GRAY,
        align=PP_ALIGN.CENTER
    )
    add_textbox(slide8, arrow,
        x + 1.0, 2.65, 0.5, 0.45,
        font_size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER
    )
    add_textbox(slide8, after,
        x + 1.5, 2.65, 1.2, 0.45,
        font_size=12, bold=True,
        color=color, align=PP_ALIGN.CENTER
    )

# Bar Chart Visual (manual bars)
add_textbox(slide8,
    "📈  Answer Relevance Comparison",
    0.4, 4.3, 6.0, 0.5,
    font_size=16, bold=True,
    color=ACCENT_CYAN
)

# Vanilla LLM Bar
vanilla_label = slide8.shapes.add_shape(
    1, Inches(0.4), Inches(5.0),
    Inches(2.5), Inches(0.4)
)
vanilla_label.fill.solid()
vanilla_label.fill.fore_color.rgb = RGBColor(0x8B, 0x00, 0x00)
vanilla_label.line.fill.background()
add_textbox(slide8, "Vanilla LLM  ~45%",
    0.5, 5.05, 2.4, 0.3,
    font_size=11, bold=True,
    color=WHITE
)

# RAG Pipeline Bar
rag_bar = slide8.shapes.add_shape(
    1, Inches(0.4), Inches(5.6),
    Inches(5.0), Inches(0.4)
)
rag_bar.fill.solid()
rag_bar.fill.fore_color.rgb = ACCENT_GREEN
rag_bar.line.fill.background()
add_textbox(slide8, "RAG Pipeline  ~90% ✅",
    0.5, 5.65, 4.8, 0.3,
    font_size=11, bold=True,
    color=DARK_BG
)

# Stats Box
add_card(slide8, 6.8, 4.3, 6.1, 2.8,
         color=RGBColor(0x06, 0x2A, 0x1E))
stats = [
    ("📚  Textbooks Processed",  "6 Chemical Engg. Books"),
    ("✂️   Total Chunks Created", "1,666 Text Chunks"),
    ("🗄️   Vector Database",      "ChromaDB"),
    ("🔍  Retrieval Strategy",    "Hybrid (Semantic + Web)"),
    ("📄  Output Format",         "PDF / DOCX Export"),
    ("🖥️   Interface",             "Streamlit UI"),
]
for i, (label, value) in enumerate(stats):
    add_textbox(slide8, label,
        6.9, 4.4 + i * 0.42, 3.0, 0.38,
        font_size=11, color=LIGHT_GRAY
    )
    add_textbox(slide8, value,
        9.9, 4.4 + i * 0.42, 2.8, 0.38,
        font_size=11, bold=True,
        color=ACCENT_GREEN
    )

# ════════════════════════════════════════════════════════
# SLIDE 9 — Tech Stack
# ════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
add_bg(slide9, prs)

add_textbox(slide9,
    "RAG Pipeline – Full Tech Stack",
    0.5, 0.3, 12.0, 0.7,
    font_size=32, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide9, 1.1, prs)

tech_stack = [
    ("🐍", "Python",          "Core Language",           RGBColor(0x1B, 0x35, 0x5C)),
    ("🦜", "LangChain",       "Pipeline Orchestration",  RGBColor(0x1B, 0x3A, 0x3A)),
    ("🗄️", "ChromaDB",        "Vector Database",         RGBColor(0x2A, 0x1B, 0x3A)),
    ("🔢", "Embeddings",      "Semantic Understanding",  RGBColor(0x1B, 0x35, 0x5C)),
    ("🔍", "Semantic Search", "Context Retrieval",       RGBColor(0x1B, 0x3A, 0x3A)),
    ("🤖", "Gemini / Groq",   "LLM Generation",          RGBColor(0x2A, 0x1B, 0x3A)),
    ("🌐", "Streamlit",       "User Interface",          RGBColor(0x1B, 0x35, 0x5C)),
    ("📄", "PDF/DOCX Export", "Output Delivery",         RGBColor(0x1B, 0x3A, 0x3A)),
    ("⚡", "FastAPI",         "REST API Backend",        RGBColor(0x2A, 0x1B, 0x3A)),
]

cols = 3
for i, (icon, name, desc, bg) in enumerate(tech_stack):
    row = i // cols
    col = i % cols
    x = 0.5 + col * 4.2
    y = 1.4 + row * 1.8
    add_card(slide9, x, y, 3.9, 1.55, color=bg)
    add_textbox(slide9, icon,
        x + 0.15, y + 0.1,
        0.8, 0.8,
        font_size=28
    )
    add_textbox(slide9, name,
        x + 1.1, y + 0.15,
        2.7, 0.5,
        font_size=16, bold=True,
        color=ACCENT_CYAN
    )
    add_textbox(slide9, desc,
        x + 1.1, y + 0.7,
        2.7, 0.5,
        font_size=12,
        color=LIGHT_GRAY
    )

# ════════════════════════════════════════════════════════
# SLIDE 10 — Key Takeaways & Conclusion
# ════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank_layout)
add_bg(slide10, prs)

add_textbox(slide10,
    "Why RAG is the Future of AI Applications",
    0.5, 0.3, 12.0, 0.7,
    font_size=30, bold=True,
    color=ACCENT_CYAN
)
add_divider(slide10, 1.1, prs)

takeaways = [
    ("✅", "No Hallucinations",
     "Responses grounded in real documents & verified sources"),
    ("✅", "No Retraining Needed",
     "Simply update the knowledge base — no expensive fine-tuning"),
    ("✅", "Domain Specific",
     "Works perfectly on private & proprietary documents"),
    ("✅", "Cost Effective",
     "Far cheaper than fine-tuning large language models"),
    ("✅", "Production Ready",
     "Scalable architecture deployable via FastAPI & Docker"),
]
for i, (icon, title, desc) in enumerate(takeaways):
    add_card(slide10, 0.4, 1.3 + i * 1.0, 12.5, 0.85)
    add_textbox(slide10, icon,
        0.55, 1.38 + i * 1.0,
        0.5, 0.6,
        font_size=18, bold=True,
        color=ACCENT_GREEN
    )
    add_textbox(slide10, title,
        1.1, 1.38 + i * 1.0,
        3.0, 0.4,
        font_size=15, bold=True,
        color=ACCENT_CYAN
    )
    add_textbox(slide10, desc,
        1.1, 1.72 + i * 1.0,
        11.5, 0.35,
        font_size=12,
        color=LIGHT_GRAY
    )

# Footer Card
add_card(slide10, 0.4, 6.4, 12.5, 0.9,
         color=RGBColor(0x06, 0x2A, 0x1E))
add_textbox(slide10,
    "🚀  Built & Deployed by Vineet Yadav",
    0.6, 6.5, 5.0, 0.4,
    font_size=14, bold=True,
    color=ACCENT_GREEN
)
add_textbox(slide10,
    "LangChain  |  ChromaDB  |  Streamlit  |  FastAPI  |  Gemini  |  Groq",
    5.5, 6.5, 7.2, 0.4,
    font_size=13,
    color=WHITE,
    align=PP_ALIGN.CENTER
)

# ════════════════════════════════════════════════════════
# SAVE PRESENTATION
# ════════════════════════════════════════════════════════
output_path = "Vineet_Yadav_RAG_Pipeline_Architecture.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"📊 Total Slides: {len(prs.slides)}")
print("🎨 Slides Generated:")
slides_list = [
    "Slide 1  →  Title Slide",
    "Slide 2  →  Problem Statement",
    "Slide 3  →  What is RAG?",
    "Slide 4  →  Full Pipeline Architecture",
    "Slide 5  →  Data Ingestion Layer",
    "Slide 6  →  Hybrid Retrieval System",
    "Slide 7  →  LLM Generation Layer",
    "Slide 8  →  Real Project Results",
    "Slide 9  →  Full Tech Stack",
    "Slide 10 →  Key Takeaways & Conclusion",
]
for s in slides_list:
    print(f"   ✅ {s}")




            